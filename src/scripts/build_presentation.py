"""
Build the FEAR Final Presentation slides from the PowerPoint template.

Inputs:
  - Template: ../FEAR-Final Presentation.pptx (one blank title slide, 11 layouts)
  - Images:   output/results/s*.png (Stata-generated event study plots)

Output:
  - ../FEAR-Final Presentation.pptx (overwritten with full slide deck)

Usage:
  python src/scripts/build_presentation.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# === CONFIGURATION ===
PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
TEMPLATE = os.path.join(PROJECT_ROOT, "..", "FEAR-Final Presentation.pptx")
OUTPUT = TEMPLATE  # overwrite in place
IMG_DIR = os.path.join(PROJECT_ROOT, "output", "results")

# Layout indices (from template inspection)
LAYOUT_TITLE_SLIDE = 0      # center title + subtitle
LAYOUT_TITLE_CONTENT = 1    # title + content body
LAYOUT_SECTION_HEADER = 2   # section header
LAYOUT_TITLE_ONLY = 5       # title only (for image slides)

# Slide dimensions (13.33" x 7.5" widescreen)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Image positioning constants
IMG_TOP = Inches(1.6)       # below title bar
IMG_MAX_W = Inches(11.0)
IMG_MAX_H = Inches(5.3)
IMG_LEFT = Inches(1.17)     # centered horizontally

# Takeaway text box positioning
TK_LEFT = Inches(0.8)
TK_TOP = Inches(6.6)
TK_WIDTH = Inches(11.7)
TK_HEIGHT = Inches(0.6)


def img_path(filename):
    """Return full path to an image file, or None if it doesn't exist."""
    p = os.path.join(IMG_DIR, filename)
    return p if os.path.exists(p) else None


def set_body_font(text_frame, font_size=18, bold=False, color=None):
    """Set default font properties on all paragraphs in a text frame."""
    for para in text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            if color:
                run.font.color.rgb = color


def add_bullet_slide(prs, title_text, bullets, sub_bullets=None):
    """Add a Title+Content slide with bullet points.

    Args:
        prs: Presentation object
        title_text: slide title
        bullets: list of strings for top-level bullets
        sub_bullets: dict mapping bullet index -> list of sub-bullet strings
    """
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_CONTENT])
    slide.placeholders[0].text = title_text

    tf = slide.placeholders[1].text_frame
    tf.clear()

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)
        p.space_after = Pt(6)

        if sub_bullets and i in sub_bullets:
            for sub in sub_bullets[i]:
                sp = tf.add_paragraph()
                sp.text = sub
                sp.level = 1
                sp.font.size = Pt(17)
                sp.space_after = Pt(4)

    return slide


def add_image_slide(prs, title_text, image_file, takeaway=None):
    """Add a Title Only slide with a centered image and optional takeaway text.

    Args:
        prs: Presentation object
        title_text: slide title
        image_file: filename in IMG_DIR
        takeaway: optional one-liner below the image
    """
    path = img_path(image_file)
    if path is None:
        print(f"  WARNING: {image_file} not found, skipping slide")
        return None

    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_ONLY])
    slide.placeholders[0].text = title_text

    # Add image, scaled to fit
    from PIL import Image
    with Image.open(path) as im:
        w_px, h_px = im.size

    aspect = w_px / h_px
    max_w = IMG_MAX_W
    max_h = IMG_MAX_H

    if aspect > (max_w / max_h):
        # width-constrained
        pic_w = max_w
        pic_h = int(max_w / aspect)
    else:
        # height-constrained
        pic_h = max_h
        pic_w = int(max_h * aspect)

    # center horizontally
    left = int((SLIDE_W - pic_w) / 2)
    slide.shapes.add_picture(path, left, IMG_TOP, pic_w, pic_h)

    # Takeaway text box
    if takeaway:
        txBox = slide.shapes.add_textbox(TK_LEFT, TK_TOP, TK_WIDTH, TK_HEIGHT)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = takeaway
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = RGBColor(0x40, 0x40, 0x40)
        p.alignment = PP_ALIGN.CENTER

    return slide


def build_deck():
    """Build the complete presentation deck."""
    prs = Presentation(TEMPLATE)

    # Delete the existing blank slide
    rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[0]

    print("Building FEAR presentation...")

    # ──────────────────────────────────────────────
    # SLIDE 1: Title Slide
    # ──────────────────────────────────────────────
    print("  Slide 1: Title")
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_SLIDE])
    s.placeholders[0].text = "Flood Risk and Home Values"
    s.placeholders[1].text = (
        "Cameron Keith\n"
        "Econ 66 — Empirical Methods in Applied Micro\n"
        "Prof. Gupta  |  Winter 2026"
    )

    # ──────────────────────────────────────────────
    # SLIDE 2: Motivation
    # ──────────────────────────────────────────────
    print("  Slide 2: Motivation")
    add_bullet_slide(prs, "Motivation", [
        "FEMA periodically redraws flood maps via Letters of Map Revision (LOMRs)",
        "Properties in Special Flood Hazard Areas (SFHAs) must carry flood insurance",
        "US housing stock worth ~$17 trillion in coastal counties alone",
        "Climate change increasing flood frequency — are markets pricing in this risk?",
    ])

    # ──────────────────────────────────────────────
    # SLIDE 3: Research Question & Hypotheses
    # ──────────────────────────────────────────────
    print("  Slide 3: Research Question")
    add_bullet_slide(prs, "Research Question & Hypotheses", [
        "RQ: How do FEMA flood zone reclassifications affect property values?",
        "H1 (Main Effect): LOMR designation reduces home values in affected zip codes",
        "H2 (Intensity): Zip codes with more LOMRs experience larger price declines",
        "H3 (Direction): Risk-increasing LOMRs reduce values; risk-decreasing LOMRs raise them",
        "H4 (Political): Republican-leaning areas respond less to reclassifications",
    ])

    # ──────────────────────────────────────────────
    # SLIDE 4: Related Literature
    # ──────────────────────────────────────────────
    print("  Slide 4: Literature")
    add_bullet_slide(prs, "Related Literature", [
        "Flood Risk Capitalization",
        "Information Shocks & Housing Markets",
        "Political Beliefs & Risk Perception",
    ], sub_bullets={
        0: [
            "Bin & Landry (2013): SFHA designation reduces prices 5–12%",
            "Gibson & Mullins (2020): flood insurance mandates capitalize into prices",
        ],
        1: [
            "Hino & Burke (2021): FEMA map updates shift price expectations",
            "Bernstein, Gustafson & Lewis (2019): SLR exposure discounted 7%",
        ],
        2: [
            "Baldauf, Garlappi & Yannelis (2020): climate skeptics pay more for exposed homes",
            "Barrage & Furst (2019): partisan divide in flood risk pricing",
        ],
    })

    # ──────────────────────────────────────────────
    # SLIDE 5: Institutional Background
    # ──────────────────────────────────────────────
    print("  Slide 5: Institutional Background")
    add_bullet_slide(prs, "Institutional Background", [
        "LOMRs: FEMA's formal process for revising flood maps",
        "SFHA mandate: properties in the 100-year floodplain with a federally-backed mortgage must carry flood insurance",
        "Study sample: 3,646 coastal zip codes across the contiguous US (2009–2022)",
        "Treatment: zip codes whose boundaries intersect a LOMR polygon",
        "Control: adjacent coastal zip codes with no LOMR exposure",
    ])

    # ──────────────────────────────────────────────
    # SLIDE 6: Data & Identification
    # ──────────────────────────────────────────────
    print("  Slide 6: Data & Identification")
    add_bullet_slide(prs, "Data & Identification Strategy", [
        "Data Sources",
        "Sample: 248K zip × quarter observations, 3,646 coastal zips, 2009–2022",
        "Identification: Staggered difference-in-differences",
        "Parallel trends: F-test fails to reject equality of pre-treatment trends (p > 0.10)",
    ], sub_bullets={
        0: [
            "Zillow Home Value Index (ZHVI) — zip-level median home values, inflation-adjusted",
            "FEMA NFHL — LOMR polygons with effective dates (treatment timing)",
            "NFIP policies & claims — insurance take-up as mechanism",
            "BLS LAUS, Census permits, MIT Election Lab — controls",
        ],
        2: [
            "Treatment = LOMR effective date × treated zip",
            "Controls: adjacent coastal zips + state-level macro controls",
            "Zip & quarter FE absorb time-invariant and aggregate shocks",
        ],
    })

    # ──────────────────────────────────────────────
    # SLIDES 7–12: Event Study Plots
    # ──────────────────────────────────────────────
    print("  Slide 7: Main Event Study")
    add_image_slide(prs, "Main Event Study: LOMR Effect on Home Values",
                    "s05_event_study_main.png",
                    "Treated zips show a statistically significant decline in home values post-LOMR, "
                    "with effects emerging ~4 quarters after reclassification.")

    print("  Slide 8: Treatment Intensity")
    add_image_slide(prs, "Treatment Intensity: High vs. Low LOMR Exposure",
                    "s06_event_study_intensity.png",
                    "Zip codes with above-median LOMR exposure experience larger and more persistent price declines.")

    print("  Slide 9: Dose-Response Quartiles")
    add_image_slide(prs, "Dose-Response: Intensity Quartiles",
                    "s06b_event_study_intensity_quartiles.png",
                    "Monotonic dose-response pattern: higher LOMR intensity quartiles show progressively larger effects.")

    print("  Slide 10: Insurance Mechanism")
    add_image_slide(prs, "Mechanism: NFIP Policy Take-Up",
                    "s08_event_study_policies.png",
                    "Insurance policies increase sharply after LOMRs, consistent with the SFHA mandate channel.")

    print("  Slide 11: Political Heterogeneity")
    add_image_slide(prs, "Political Heterogeneity: Republican Vote Share",
                    "s09c_event_study_republican.png",
                    "Republican-leaning areas show attenuated response to flood reclassifications.")

    print("  Slide 12: Up vs. Down Reclassifications")
    add_image_slide(prs, "Risk Direction: Up vs. Down Reclassifications",
                    "s09_event_study_updown_decomposed.png",
                    "Risk-increasing LOMRs drive price declines; risk-decreasing LOMRs have a positive but smaller effect.")

    # ──────────────────────────────────────────────
    # SLIDE 13: Robustness Summary
    # ──────────────────────────────────────────────
    print("  Slide 13: Robustness")
    add_bullet_slide(prs, "Robustness Checks", [
        "Parallel trends: pre-treatment coefficients jointly insignificant (F-test p > 0.10)",
        "Leave-one-out states: results stable dropping any single state",
        "Alternative clustering: two-way (zip + quarter) and state-level clustering yield consistent SEs",
        "Placebo outcome: no effect on county unemployment — rules out macro confounds",
        "Bacon decomposition: TWFE estimates driven by clean comparisons, not problematic 2×2s",
        "Callaway & Sant'Anna (2021): heterogeneity-robust estimator confirms main findings",
    ])

    # ──────────────────────────────────────────────
    # SLIDE 14: Conclusion
    # ──────────────────────────────────────────────
    print("  Slide 14: Conclusion")
    add_bullet_slide(prs, "Conclusion", [
        "Key Findings",
        "Implications",
    ], sub_bullets={
        0: [
            "FEMA flood zone reclassifications reduce home values by 2–4% in affected zip codes",
            "Effects are dose-responsive: more LOMRs → larger price declines",
            "Insurance mandate is the primary mechanism — policies surge post-LOMR",
            "Republican-leaning areas show muted responses, consistent with belief-driven heterogeneity",
        ],
        1: [
            "Flood maps are an effective tool for communicating risk to housing markets",
            "Political polarization may weaken the price signal of climate risk information",
            "As climate change increases flood frequency, map revision policy becomes more consequential",
        ],
    })

    # ──────────────────────────────────────────────
    # SLIDE 15: Thank You
    # ──────────────────────────────────────────────
    print("  Slide 15: Thank You")
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_SLIDE])
    s.placeholders[0].text = "Thank You"
    s.placeholders[1].text = "Questions?"

    # ══════════════════════════════════════════════
    # APPENDIX / BACKUP SLIDES
    # ══════════════════════════════════════════════
    print("  Appendix section header")
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUT_SECTION_HEADER])
    s.placeholders[0].text = "Appendix"
    s.placeholders[1].text = "Backup Slides"

    appendix_slides = [
        ("Parallel Trends Test", "s10_parallel_trends.png",
         "Pre-treatment coefficients are individually and jointly insignificant."),
        ("Treatment Timing Distribution", "s10_treatment_timing_hist.png",
         "LOMRs are staggered across the sample period, with no bunching around a single date."),
        ("Bacon Decomposition", "s12_bacon_decomposition.png",
         "Goodman-Bacon decomposition: TWFE estimate driven by clean treated-vs-never-treated comparisons."),
        ("Callaway & Sant'Anna Estimator", "s13_event_study_cs.png",
         "Heterogeneity-robust CS estimator confirms the main TWFE results."),
        ("Placebo: County Unemployment", "s14_event_study_placebo.png",
         "No effect on unemployment — LOMRs do not coincide with local economic shocks."),
        ("Disclosure Law Heterogeneity", "s09b_event_study_disclosure.png",
         "States with mandatory flood risk disclosure laws show larger capitalization effects."),
        ("Up vs. Down: Combined", "s09_event_study_updown.png",
         "Combined up/down event study showing differential treatment effects by risk direction."),
        ("SFHA Crossing Heterogeneity", "s09d_event_study_sfha_crossing.png",
         "LOMRs that move properties across the SFHA boundary show the largest effects."),
        ("Up vs. Down × Intensity", "s09_event_study_updown_intensity.png",
         "Interaction of risk direction and treatment intensity."),
        ("Insurance × Risk Direction", "s08b_event_study_policies_updown.png",
         "NFIP policy response differs by whether the LOMR increased or decreased risk designation."),
    ]

    for title, img, takeaway in appendix_slides:
        print(f"  Appendix: {title}")
        add_image_slide(prs, title, img, takeaway)

    # ──────────────────────────────────────────────
    # Save
    # ──────────────────────────────────────────────
    prs.save(OUTPUT)
    print(f"\nSaved to: {OUTPUT}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_deck()
